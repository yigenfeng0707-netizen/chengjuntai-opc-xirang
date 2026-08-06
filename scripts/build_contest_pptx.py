# -*- coding: utf-8 -*-
"""生成息壤杯预赛提交用专业答辩 PPTX（图文并茂 · ≤100MB）。

用法:
  python scripts/render_ppt_assets.py   # 生成 assets/ppt/*.png
  python scripts/build_contest_pptx.py  # 嵌入并写出 docs/成军台_息壤杯预赛答辩.pptx
  # 或一步：python scripts/build_contest_pptx.py --render
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets" / "ppt"
OUT = ROOT / "docs" / "成军台_息壤杯预赛答辩.pptx"

INK = RGBColor(0x0F, 0x1C, 0x2E)
TEAL = RGBColor(0x0D, 0x94, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# (asset filename, speaker notes)
SLIDES: list[tuple[str, str]] = [
    (
        "01_cover.png",
        "各位评委好，我们是成军台。一句话：让一个人也能成军——"
        "在息壤上把调研、内容、问数、运营、复盘拉成一支可验收的 AI 小队；"
        "同时接上真政采与标书，服务政企落地。",
    ),
    (
        "02_pain.png",
        "OPC 不是缺模型，是缺操作系统：目标进来，任务树、人审、产物和周报要能交付。"
        "电信政企场景里，还要把种草内容、标讯洞察和标书材料串在同一条链上。",
    ),
    (
        "03_solution.png",
        "方案：目标驱动战役 + AI 员工矩阵 + 可导出 Word。"
        "增量：内容工厂、标书工作台、真政采 NL2SQL。"
        "底座落在息壤与天翼云，赛道是惠民 AI+自选。",
    ),
    (
        "04_timeline_60s.png",
        "零讲解路径：点评委体验 → 打开已完成样例 → 导出成军周报 → 再看一张问数表。"
        "评委自己点也能走通——直接打实用性与完整性。",
    ),
    (
        "05_agent_matrix.png",
        "一人输入目标，五类角色并行，产物进抽屉，可预览、可导出 Word——"
        "创新点在可验收小队，不在多聊几句。",
    ),
    (
        "06_architecture.png",
        "工程底座：Web 看板、司令 Agent、战役状态机、多 Agent、"
        "息壤 LLM、NL2SQL 与标讯库。无 Key 明确失败，禁止静默 mock。",
    ),
    (
        "07_llm_provider.png",
        "双轨就绪：竞赛 Token 到手切 primary；之前用过渡模型也能真实 E2E。"
        "没有 Key 绝不假装成功——这是技术诚信，也是评委信任。",
    ),
    (
        "08_metrics.png",
        "312条浙江政采真实标讯入库 PostgreSQL；15个 MCP 工具；"
        "样例周包种子化保证 60 秒可走通；产物可预览、可导出 Word。",
    ),
    (
        "09_evidence_matrix.png",
        "标书材料工作台：粘贴招标文本 → 要求拆解 → 生成证据矩阵 → 导出 Word 附卷。"
        "电信员工熟悉的政企场景，OPC 也能成军。",
    ),
    (
        "10_five_dimensions.png",
        "五维对位：创新在 OS 形态，实用在周报与标讯成军，完整在 60 秒可走通，"
        "技术在多智能体与真数据，价值在惠民一人成军与电信政企落地。",
    ),
    (
        "11_business_value.png",
        "先把一人成军跑通，再按 Token 与席位经营；模板与标讯越用越厚。"
        "惠民：让个体用得起、看得见交付的 AI 协作。",
    ),
    (
        "12_demo_proof.png",
        "提交对齐：60 秒电影级视频、本 PPT、主仓 chengjuntai-opc-xirang、"
        "评委只读账号当面提供；公网 HTTPS 与正式 Token 按审批推进。",
    ),
    (
        "13_roadmap.png",
        "预赛打穿 60 秒与材料完整；复赛加深行业模板、证据矩阵与经营数据；"
        "远景行业方案包与孵化。",
    ),
    (
        "14_thanks.png",
        "谢谢各位评委。欢迎提问，也可现在打开 Demo，点「评委 60 秒体验」自己走一遍。",
    ),
]


def _blank(prs: Presentation):
    return prs.slide_layouts[6]


def _add_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def _footer_chip(slide, page: int, total: int) -> None:
    """Subtle page marker (assets already brand-headed; keep tiny for projection)."""
    box = slide.shapes.add_textbox(
        Inches(12.0), Inches(7.15), Inches(1.1), Inches(0.28)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page}/{total}"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED
    run.font.name = "Microsoft YaHei"


def ensure_assets(do_render: bool) -> None:
    missing = [name for name, _ in SLIDES if not (ASSETS / name).exists()]
    if missing or do_render:
        sys.path.insert(0, str(ROOT / "scripts"))
        from render_ppt_assets import render_all  # noqa: WPS433

        render_all()
    still = [name for name, _ in SLIDES if not (ASSETS / name).exists()]
    if still:
        raise FileNotFoundError(
            f"Missing PPT assets: {still}. Run scripts/render_ppt_assets.py"
        )


def build(do_render: bool = False) -> Path:
    ensure_assets(do_render)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    total = len(SLIDES)

    for i, (name, notes) in enumerate(SLIDES, start=1):
        path = ASSETS / name
        s = prs.slides.add_slide(_blank(prs))
        # full-bleed visual
        s.shapes.add_picture(str(path), Emu(0), Emu(0), width=SLIDE_W, height=SLIDE_H)
        _footer_chip(s, i, total)
        _add_notes(s, notes)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    size_mb = OUT.stat().st_size / (1024 * 1024)
    print(f"OK {OUT}")
    print(f"slides={total}  size={size_mb:.2f} MB  assets={ASSETS}")
    return OUT


def main():
    ap = argparse.ArgumentParser(description="Build 成军台 contest PPTX")
    ap.add_argument(
        "--render", action="store_true", help="Re-render Pillow assets first"
    )
    args = ap.parse_args()
    build(do_render=args.render)


if __name__ == "__main__":
    main()
